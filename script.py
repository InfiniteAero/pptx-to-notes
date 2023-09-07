# tried to make this as clean as possible but it seems wrong somehow
# TODO: GUI and user-friendliness

from pptx import Presentation
from docx import Document
# import argparse
import sys

# parser = argparse.ArgumentParser(
#                     prog='pptx_to_notes',
#                     description='Converts a pptx presentation into a abridged word doc (i.e. notes)',
#                     epilog='very cool')


def init_presentation(path):
    """Initializes a Presentation object representing a pptx file"""
    try:
        prs = Presentation(path)
    except FileNotFoundError:
        print(
            "No file exists at that file location. Double check your file path to make sure it's correct."
        )
        sys.exit(1)
    return prs


def extract_text(prs):
    """Extracts all the text from the given presentation for later"""
    extracted_text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    extracted_text.append(run.text)
    return extracted_text


def create_document(extracted_text):
    """Creates a word document with the extracted slideshow text"""
    document = Document()
    for run in extracted_text:
        document.add_paragraph(run)
    return document


# def main():
#     parser.add_argument('file_path')
#     parser.add_argument('save_path') # needs to be an absolute file path
# 
#     args = parser.parse_args()
# 
#     prs = init_presentation(args.file_path) # this will definitely cause issues lol
#     extracted_text = extract_text(prs)
#     document = create_document(extracted_text)
#     document.save(args.save_path)
# 
#     print("Done!")
# 
# 
# if __name__ == "__main__":
#     main()
# 